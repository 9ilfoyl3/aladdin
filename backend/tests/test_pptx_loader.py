"""测试 PptxLoader（PPT 文档加载器）"""

import os

import pytest
from pptx import Presentation

from app.pipeline.loaders.pptx_loader import PptxLoader
from app.pipeline.loader import LoadResult


class TestPptxLoader:
    """PptxLoader 单元测试"""

    def setup_method(self):
        self.loader = PptxLoader()

    def _create_pptx(self, tmp_path, filename: str, slides_content: list[list[str]]) -> str:
        """辅助方法：创建包含指定幻灯片内容的 pptx 文件

        Args:
            slides_content: 每个元素是一张幻灯片的文本列表
        """
        file_path = str(tmp_path / filename)
        prs = Presentation()
        for texts in slides_content:
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # 空白布局
            for text in texts:
                txBox = slide.shapes.add_textbox(
                    left=0, top=0, width=prs.slide_width, height=prs.slide_height
                )
                txBox.text_frame.paragraphs[0].text = text
            prs.save(file_path)
        return file_path

    def test_load_single_slide(self, tmp_path):
        """加载单幻灯片 pptx"""
        file_path = self._create_pptx(tmp_path, "single.pptx", [["Hello PPT"]])

        result = self.loader.load(file_path)

        assert isinstance(result, LoadResult)
        assert "Hello PPT" in result.content
        assert "[Slide 1]" in result.content
        assert result.metadata["filename"] == "single.pptx"
        assert result.metadata["file_type"] == "pptx"
        assert result.metadata["file_size"] > 0
        assert result.metadata["slide_count"] == 1

    def test_load_multi_slide(self, tmp_path):
        """加载多幻灯片 pptx"""
        slides = [["第一页内容"], ["第二页内容"], ["第三页内容"]]
        file_path = self._create_pptx(tmp_path, "multi.pptx", slides)

        result = self.loader.load(file_path)

        assert "[Slide 1]" in result.content
        assert "[Slide 2]" in result.content
        assert "[Slide 3]" in result.content
        assert "第一页内容" in result.content
        assert "第二页内容" in result.content
        assert "第三页内容" in result.content
        assert result.metadata["slide_count"] == 3

    def test_slides_separated_by_double_newline(self, tmp_path):
        """幻灯片之间用双换行分隔"""
        slides = [["幻灯片A"], ["幻灯片B"]]
        file_path = self._create_pptx(tmp_path, "sep.pptx", slides)

        result = self.loader.load(file_path)

        assert "\n\n" in result.content

    def test_file_not_found(self):
        """文件不存在时抛出 FileNotFoundError"""
        with pytest.raises(FileNotFoundError, match="文件不存在"):
            self.loader.load("/nonexistent/path/file.pptx")

    def test_invalid_pptx(self, tmp_path):
        """无效 pptx 文件抛出 ValueError"""
        file_path = tmp_path / "invalid.pptx"
        file_path.write_text("this is not a pptx", encoding="utf-8")

        with pytest.raises(ValueError, match="无法解析 pptx 文件"):
            self.loader.load(str(file_path))

    def test_file_size_metadata(self, tmp_path):
        """file_size 反映实际文件字节数"""
        file_path = self._create_pptx(tmp_path, "size.pptx", [["测试内容"]])

        result = self.loader.load(file_path)

        actual_size = os.path.getsize(file_path)
        assert result.metadata["file_size"] == actual_size

    def test_multiple_shapes_per_slide(self, tmp_path):
        """单张幻灯片包含多个文本框"""
        slides = [["文本框1", "文本框2"]]
        file_path = self._create_pptx(tmp_path, "multi_shape.pptx", slides)

        result = self.loader.load(file_path)

        assert "文本框1" in result.content
        assert "文本框2" in result.content
