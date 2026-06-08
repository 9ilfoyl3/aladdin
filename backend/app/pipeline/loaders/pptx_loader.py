"""PPT (pptx) 文档加载器（基于 python-pptx）

支持提取文本和嵌入图片。图片写入临时目录（避免内存压力），
由 pipeline 的 OCR 流程处理后清理。
"""

import hashlib
import os
import tempfile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.pipeline.loader import BaseLoader, EmbeddedImage, LoadResult


# 最小图片数据大小（字节），过小的图片跳过
_MIN_IMAGE_BYTES = 1024
# 单文档最大提取图片数量
_MAX_IMAGES_PER_DOC = 50


class PptxLoader(BaseLoader):
    """处理 .pptx 文件的加载器，同时提取文本和嵌入图片"""

    def load(self, file_path: str) -> LoadResult:
        """加载 PPT 文件，提取全部幻灯片文本和嵌入图片

        格式：每张幻灯片以 "[Slide N]" 标题开头，幻灯片之间以双换行分隔。
        图片写入临时目录，通过 content_hash 去重。

        Args:
            file_path: 文件路径

        Returns:
            LoadResult: 包含文件内容、按页文本、元数据和嵌入图片列表

        Raises:
            FileNotFoundError: 文件不存在
            ValueError: 无效的 pptx 文件
        """
        if not os.path.isfile(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")

        file_size = os.path.getsize(file_path)

        try:
            prs = Presentation(file_path)
        except Exception as e:
            raise ValueError(f"无法解析 pptx 文件: {file_path}，错误: {e}")

        tmp_dir = tempfile.mkdtemp(prefix="pptx_images_")

        slides_text: list[str] = []
        images: list[EmbeddedImage] = []
        seen_hashes: set[str] = set()

        for idx, slide in enumerate(prs.slides, start=1):
            paragraphs = []
            for shape in slide.shapes:
                # 提取文本
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            paragraphs.append(text)

                # 提取图片
                if (
                    shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                    and len(images) < _MAX_IMAGES_PER_DOC
                ):
                    self._extract_shape_image(
                        shape, idx, tmp_dir, images, seen_hashes
                    )

            slide_content = f"[Slide {idx}]\n" + "\n".join(paragraphs)
            slides_text.append(slide_content)

        slide_count = len(slides_text)
        content = "\n\n".join(slides_text)

        metadata = {
            "filename": os.path.basename(file_path),
            "file_type": "pptx",
            "file_size": file_size,
            "slide_count": slide_count,
            "embedded_image_count": len(images),
        }

        return LoadResult(
            content=content,
            metadata=metadata,
            images=images,
            page_texts=slides_text,
        )

    @staticmethod
    def _extract_shape_image(
        shape,
        slide_idx: int,
        tmp_dir: str,
        images: list[EmbeddedImage],
        seen_hashes: set[str],
    ) -> None:
        """从 shape 中提取图片并写入临时目录

        Args:
            shape: pptx shape 对象
            slide_idx: 幻灯片序号（从1开始）
            tmp_dir: 临时目录路径
            images: 图片列表（会被修改）
            seen_hashes: 已见 hash 集合（会被修改）
        """
        try:
            image = shape.image
            image_bytes = image.blob
        except Exception:
            return

        if not image_bytes or len(image_bytes) < _MIN_IMAGE_BYTES:
            return

        # 计算 hash 去重
        img_hash = hashlib.md5(image_bytes).hexdigest()
        if img_hash in seen_hashes:
            return
        seen_hashes.add(img_hash)

        # 从 content_type 推断格式
        content_type = image.content_type or ""
        img_format = "png"
        if "jpeg" in content_type or "jpg" in content_type:
            img_format = "jpeg"
        elif "png" in content_type:
            img_format = "png"

        img_filename = f"slide{slide_idx}_img{len(images)+1}_{img_hash[:8]}.{img_format}"
        img_path = os.path.join(tmp_dir, img_filename)

        with open(img_path, "wb") as f:
            f.write(image_bytes)

        images.append(
            EmbeddedImage(
                file_path=img_path,
                format=img_format,
                page_or_index=slide_idx,
                content_hash=img_hash,
                description=f"slide{slide_idx}_img{len(images)+1}",
            )
        )
